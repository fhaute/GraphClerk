from __future__ import annotations

from io import BytesIO
from typing import Any

from app.core.config import Settings
from app.models.artifact import Artifact
from app.models.enums import Modality, SourceFidelity
from app.schemas.evidence_unit_candidate import EvidenceUnitCandidate
from app.services.errors import (
    ExtractionReturnedNoEvidenceError,
    ExtractorUnavailableError,
    PptxExtractionError,
)
from app.services.raw_source_store import RawSourceStore

try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore[misc, assignment]


def _flatten_table_text(shape: Any) -> str:
    rows: list[str] = []
    for row in shape.table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append("\t".join(cells))
    return "\n".join(rows).strip()


class PptxExtractor:
    """Extract text from PowerPoint (.pptx) slides using python-pptx (Phase 5 Slice E).

    Requires the ``pptx`` optional extra. Does not persist or mutate ``Artifact``.
    """

    def __init__(self, *, settings: Settings) -> None:
        if Presentation is None:
            raise ExtractorUnavailableError(
                "PPTX extraction requires the optional `pptx` dependency "
                "(e.g. pip install -e '.[pptx]' for graphclerk-backend).",
            )
        self._settings = settings
        self._raw_store = RawSourceStore(settings)

    def extract(self, artifact: Artifact) -> list[EvidenceUnitCandidate]:
        if artifact.artifact_type != "pptx":
            raise PptxExtractionError(
                f"PptxExtractor expected artifact_type 'pptx', got {artifact.artifact_type!r}.",
            )

        pptx_bytes = self._raw_store.read_persisted_bytes(
            storage_uri=artifact.storage_uri,
            checksum=artifact.checksum,
            filename=artifact.filename,
            raw_text=artifact.raw_text,
        )

        try:
            prs = Presentation(BytesIO(pptx_bytes))
        except Exception as e:
            raise PptxExtractionError(f"Invalid or unreadable PPTX: {e}") from e

        slide_count = len(prs.slides)
        if slide_count == 0:
            raise ExtractionReturnedNoEvidenceError()

        candidates: list[EvidenceUnitCandidate] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            meta: dict[str, Any] = {"slide_count": slide_count, "extractor": "python-pptx"}

            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                tt = title_shape.text.strip()
                if tt:
                    candidates.append(
                        EvidenceUnitCandidate(
                            modality=Modality.slide,
                            content_type="slide_title",
                            text=tt,
                            location={"slide_number": slide_idx, "region": "title"},
                            source_fidelity=SourceFidelity.extracted,
                            confidence=1.0,
                            metadata=meta,
                        )
                    )

            for shape in slide.shapes:
                if shape.has_table:
                    try:
                        tbl = _flatten_table_text(shape)
                    except Exception as e:
                        msg = f"Failed to read table on slide {slide_idx}: {e}"
                        raise PptxExtractionError(msg) from e
                    if tbl:
                        candidates.append(
                            EvidenceUnitCandidate(
                                modality=Modality.slide,
                                content_type="slide_table_text",
                                text=tbl,
                                location={"slide_number": slide_idx, "region": "table"},
                                source_fidelity=SourceFidelity.extracted,
                                confidence=1.0,
                                metadata=meta,
                            )
                        )
                    continue

                if not shape.has_text_frame:
                    continue
                if title_shape is not None and shape is title_shape:
                    continue
                body = shape.text.strip()
                if body:
                    candidates.append(
                        EvidenceUnitCandidate(
                            modality=Modality.slide,
                            content_type="slide_body_text",
                            text=body,
                            location={"slide_number": slide_idx, "region": "body"},
                            source_fidelity=SourceFidelity.extracted,
                            confidence=1.0,
                            metadata=meta,
                        )
                    )

            if slide.has_notes_slide:
                try:
                    ns = slide.notes_slide
                    ntf = ns.notes_text_frame
                    if ntf is not None:
                        nt = ntf.text.strip()
                        if nt:
                            candidates.append(
                                EvidenceUnitCandidate(
                                    modality=Modality.slide,
                                    content_type="slide_notes_text",
                                    text=nt,
                                    location={"slide_number": slide_idx, "region": "notes"},
                                    source_fidelity=SourceFidelity.extracted,
                                    confidence=1.0,
                                    metadata=meta,
                                )
                            )
                except Exception as e:
                    msg = f"Failed to read notes on slide {slide_idx}: {e}"
                    raise PptxExtractionError(msg) from e

        if not candidates:
            raise ExtractionReturnedNoEvidenceError()

        return candidates
