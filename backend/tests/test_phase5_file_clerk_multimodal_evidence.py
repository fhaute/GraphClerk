"""Phase 5 Slice I — File Clerk + graph compatibility for multimodal EvidenceUnits.

Proves ``POST /retrieve`` yields ``RetrievalPacket.evidence_units`` entries that preserve
PDF and PPTX metadata (modality, content_type, source_fidelity, location, ids, text)
when those units are linked via existing graph node evidence APIs and reached through
the same stubbed semantic search pattern as Phase 4 retrieve tests.

Image and audio extractors currently validate inputs then respond with **503** because
OCR/ASR backends are not configured, so they **do not emit EvidenceUnits**. Packet coverage
for ``modality=image`` / ``modality=audio`` evidence rows is deferred until those
extractors produce candidates.
"""

from __future__ import annotations

import uuid

import httpx
import pytest
from httpx import ASGITransport

from app.db.session import get_sessionmaker
from app.main import create_app
from app.models.enums import SemanticIndexVectorStatus
from app.models.semantic_index import SemanticIndex
from app.services.embedding_adapter import DeterministicFakeEmbeddingAdapter
from app.services.embedding_service import EmbeddingService
from app.services.semantic_index_search_service import SemanticIndexSearchService
from app.services.vector_index_service import SearchHit

# Minimal 1-page PDF with selectable text (same fixture as ``test_phase5_pdf_extractor``).
_MINIMAL_HELLO_PDF = (
    b"%PDF-1.4\n"
    b"1 0 obj<< /Type /Catalog /Pages 2 0 R >>endobj\n"
    b"2 0 obj<< /Type /Pages /Kids [3 0 R] /Count 1 >>endobj\n"
    b"3 0 obj<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
    b"/Contents 4 0 R /Resources<< /Font<< /F1 5 0 R >> >> >>endobj\n"
    b"4 0 obj<< /Length 44 >>stream\n"
    b"BT /F1 24 Tf 100 700 Td (Hello) Tj ET\n"
    b"endstream\n"
    b"endobj\n"
    b"5 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj\n"
    b"xref\n"
    b"0 6\n"
    b"0000000000 65535 f \n"
    b"0000000009 00000 n \n"
    b"0000000058 00000 n \n"
    b"0000000115 00000 n \n"
    b"0000000266 00000 n \n"
    b"0000000355 00000 n \n"
    b"trailer<< /Size 6 /Root 1 0 R >>\n"
    b"startxref\n"
    b"433\n"
    b"%%EOF"
)


def _minimal_pptx_bytes() -> bytes:
    pytest.importorskip("pptx")
    from io import BytesIO

    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Slice E Title"
    body = slide.placeholders[1]
    body.text = "Body paragraph for tests."
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _mark_indexes_indexed(*, semantic_index_ids: list[str]) -> None:
    SessionMaker = get_sessionmaker()
    with SessionMaker() as session:
        for sid in semantic_index_ids:
            row = session.get(SemanticIndex, uuid.UUID(sid))
            assert row is not None
            row.vector_status = SemanticIndexVectorStatus.indexed
        session.commit()


class _StubVectorIndexService:
    def __init__(self, hits: list[SearchHit]) -> None:
        self._hits = hits

    def search_semantic_indexes(
        self, *, query_vector: list[float], limit: int = 5
    ) -> list[SearchHit]:
        _ = query_vector
        _ = limit
        return self._hits


def _search_service_factory(*, hits: list[SearchHit]):
    def factory(*, session):
        embedding = EmbeddingService(
            adapter=DeterministicFakeEmbeddingAdapter(dimension=8),
            expected_dimension=8,
        )
        vector = _StubVectorIndexService(hits=hits)
        return SemanticIndexSearchService(
            session=session,
            embedding_service=embedding,
            vector_index_service=vector,
        )

    return factory


@pytest.mark.asyncio
async def test_retrieve_includes_pdf_evidence_unit_when_linked_to_graph(
    db_ready: None,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    pytest.importorskip("pypdf")

    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        ing = await client.post(
            "/artifacts",
            files={"file": ("hello.pdf", _MINIMAL_HELLO_PDF, "application/pdf")},
        )
        assert ing.status_code == 200
        artifact_id = ing.json()["artifact_id"]

        ev_res = await client.get(f"/artifacts/{artifact_id}/evidence")
        assert ev_res.status_code == 200
        pdf_units = [x for x in ev_res.json()["items"] if x.get("content_type") == "pdf_page_text"]
        assert len(pdf_units) >= 1
        evidence_unit_id = pdf_units[0]["id"]
        assert pdf_units[0]["text"]
        assert pdf_units[0]["modality"] == "pdf"
        assert pdf_units[0]["source_fidelity"] == "extracted"
        assert pdf_units[0]["location"] is not None
        assert pdf_units[0]["location"].get("page") == 1

        node_resp = await client.post(
            "/graph/nodes",
            json={"node_type": "concept", "label": "PdfEvidenceNode"},
        )
        node_id = node_resp.json()["id"]
        link = await client.post(
            f"/graph/nodes/{node_id}/evidence",
            json={
                "evidence_unit_id": evidence_unit_id,
                "support_type": "supports",
                "confidence": 0.9,
            },
        )
        assert link.status_code == 200

        sid = (
            await client.post(
                "/semantic-indexes",
                json={"meaning": "pdf_multimodal_slice_i", "entry_node_ids": [node_id]},
            )
        ).json()["id"]

    _mark_indexes_indexed(semantic_index_ids=[sid])

    hits = [
        SearchHit(
            semantic_index_id=uuid.UUID(sid),
            score=0.92,
            payload={"semantic_index_id": sid},
        ),
    ]
    import app.services.semantic_index_search_factory as factory_module

    monkeypatch.setattr(
        factory_module,
        "build_semantic_index_search_service",
        _search_service_factory(hits=hits),
    )

    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post("/retrieve", json={"question": "pdf multimodal evidence probe"})
        assert res.status_code == 200
        body = res.json()
        assert body["packet_type"] == "retrieval_packet"
        units = body["evidence_units"]
        match = next((u for u in units if u["evidence_unit_id"] == evidence_unit_id), None)
        assert match is not None, f"expected evidence_unit_id {evidence_unit_id} in {units!r}"
        assert match["artifact_id"] == artifact_id
        assert match["modality"] == "pdf"
        assert match["content_type"] == "pdf_page_text"
        assert match["source_fidelity"] == "extracted"
        assert match["location"] is not None
        assert match["location"].get("page") == 1
        assert match["text"]
        assert "Hello" in match["text"]


@pytest.mark.asyncio
async def test_retrieve_includes_pptx_evidence_unit_when_linked_to_graph(
    db_ready: None,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    pytest.importorskip("pptx")

    pptx_bytes = _minimal_pptx_bytes()
    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        ing = await client.post(
            "/artifacts",
            files={
                "file": (
                    "deck.pptx",
                    pptx_bytes,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
        assert ing.status_code == 200
        artifact_id = ing.json()["artifact_id"]

        ev_res = await client.get(f"/artifacts/{artifact_id}/evidence")
        assert ev_res.status_code == 200
        slide_units = [
            x
            for x in ev_res.json()["items"]
            if x.get("content_type") in ("slide_title", "slide_body_text")
        ]
        assert len(slide_units) >= 1
        evidence_unit_id = slide_units[0]["id"]
        assert slide_units[0]["text"]
        assert slide_units[0]["modality"] == "slide"
        assert slide_units[0]["source_fidelity"] == "extracted"
        loc = slide_units[0]["location"]
        assert loc is not None
        assert loc.get("slide_number") is not None
        assert loc.get("region")

        node_resp = await client.post(
            "/graph/nodes",
            json={"node_type": "concept", "label": "PptxEvidenceNode"},
        )
        node_id = node_resp.json()["id"]
        link = await client.post(
            f"/graph/nodes/{node_id}/evidence",
            json={
                "evidence_unit_id": evidence_unit_id,
                "support_type": "supports",
                "confidence": 0.85,
            },
        )
        assert link.status_code == 200

        sid = (
            await client.post(
                "/semantic-indexes",
                json={"meaning": "pptx_multimodal_slice_i", "entry_node_ids": [node_id]},
            )
        ).json()["id"]

    _mark_indexes_indexed(semantic_index_ids=[sid])

    hits = [
        SearchHit(
            semantic_index_id=uuid.UUID(sid),
            score=0.93,
            payload={"semantic_index_id": sid},
        ),
    ]
    import app.services.semantic_index_search_factory as factory_module

    monkeypatch.setattr(
        factory_module,
        "build_semantic_index_search_service",
        _search_service_factory(hits=hits),
    )

    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post("/retrieve", json={"question": "pptx multimodal evidence probe"})
        assert res.status_code == 200
        body = res.json()
        units = body["evidence_units"]
        match = next((u for u in units if u["evidence_unit_id"] == evidence_unit_id), None)
        assert match is not None, f"expected evidence_unit_id {evidence_unit_id} in {units!r}"
        assert match["artifact_id"] == artifact_id
        assert match["modality"] == "slide"
        assert match["content_type"] in ("slide_title", "slide_body_text")
        assert match["source_fidelity"] == "extracted"
        assert match["location"] is not None
        assert match["location"].get("slide_number") is not None
        assert match["location"].get("region")
        assert match["text"]
