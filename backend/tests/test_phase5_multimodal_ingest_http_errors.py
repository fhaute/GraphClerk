"""HTTP error matrix for ``POST /artifacts`` multimodal ingestion (Phase 5 Slice J).

Status code split (contract for docs/audit):

- **400** — ``ExtractorNotRegisteredError`` (nothing registered for that modality),
  resolver failures (unsupported type, video), ``ExtractionReturnedNoEvidenceError``,
  and other ``GraphClerkError`` subclasses such as parse errors (e.g. corrupt PDF/PPTX).

- **503** — ``ExtractorUnavailableError``: extractor **is** registered for the modality but
  cannot run (optional dependency placeholder, OCR/ASR shell, or a stub that raises this).

Bad input after a successful artifact write still maps to **400** on the response; the
route does not treat parse failures as **503**.
"""

from __future__ import annotations

import importlib
import io
import wave
from io import BytesIO

import httpx
import pytest
from httpx import ASGITransport

from app.api.routes import artifacts as artifacts_routes
from app.main import create_app
from app.models.artifact import Artifact
from app.models.enums import Modality
from app.schemas.evidence_unit_candidate import EvidenceUnitCandidate
from app.services.errors import ExtractorUnavailableError
from app.services.extraction import ExtractorRegistry

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# Minimal 1-page PDF (same bytes as ``test_phase5_pdf_extractor``).
_MINIMAL_HELLO_PDF = (
    b"%PDF-1.4\n"
    b"1 0 obj<< /Type /Catalog /Pages 2 0 R >>endobj\n"
    b"2 0 obj<< /Type /Pages /Kids [3 0 R] /Count 1 >>endobj\n"
    b"3 0 obj<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
    b"/Contents 4 0 R /Resources<< /Font<< /F1 5 0 R >> >> >>endobj\n"
    b"4 0 obj<< /Length 44 >>stream\n"
    b"BT /F1 24 Tf 100 700 Td (Hello) Tj ET\n"
    b"endstream\n"
    b"endobj\n"
    b"5 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj\n"
    b"xref\n"
    b"0 6\n"
    b"0000000000 65535 f \n"
    b"0000000009 00000 n \n"
    b"0000000058 00000 n \n"
    b"0000000115 00000 n \n"
    b"0000000266 00000 n \n"
    b"0000000355 00000 n \n"
    b"trailer<< /Size 6 /Root 1 0 R >>\n"
    b"startxref\n"
    b"433\n"
    b"%%EOF"
)

# Minimal valid 1x1 RGBA PNG (``test_phase5_image_extractor``).
_MINIMAL_PNG = (
    b"\x89PNG\r\n\x1a\n"
    b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89"
    b"\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01\r\n-\xb4"
    b"\x00\x00\x00\x00IEND\xaeB`\x82"
)


def _minimal_wav_bytes() -> bytes:
    buf = io.BytesIO()
    with wave.open(buf, "wb") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(8000)
        wf.writeframes(b"\x00\x00" * 8)
    return buf.getvalue()


def _blank_pdf_bytes() -> bytes:
    pytest.importorskip("pypdf")
    from pypdf import PdfWriter

    buf = BytesIO()
    w = PdfWriter()
    w.add_blank_page(width=72, height=72)
    w.add_blank_page(width=72, height=72)
    w.write(buf)
    return buf.getvalue()


def _minimal_pptx_bytes() -> bytes:
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Title"
    body = slide.placeholders[1]
    body.text = "Body for http error matrix."
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _empty_text_pptx_bytes() -> bytes:
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    layouts = prs.slide_layouts
    layout_idx = 6 if len(layouts) > 6 else max(0, len(layouts) - 1)
    prs.slides.add_slide(layouts[layout_idx])
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


class _UnavailableExtractor:
    def extract(self, artifact: Artifact) -> list[EvidenceUnitCandidate]:
        raise ExtractorUnavailableError("extractor stub unavailable")


@pytest.mark.asyncio
async def test_post_corrupt_pdf_returns_400_when_pypdf_available(db_ready: None) -> None:
    pytest.importorskip("pypdf")
    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post(
            "/artifacts",
            files={"file": ("bad.pdf", b"not a pdf", "application/pdf")},
        )
    assert res.status_code == 400
    assert "detail" in res.json()


@pytest.mark.asyncio
async def test_post_blank_pdf_returns_400_extraction_returned_no_evidence(db_ready: None) -> None:
    pytest.importorskip("pypdf")
    blank = _blank_pdf_bytes()
    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post(
            "/artifacts",
            files={"file": ("blank.pdf", blank, "application/pdf")},
        )
    assert res.status_code == 400
    assert "extraction_returned_no_evidence" in res.json()["detail"]


@pytest.mark.asyncio
async def test_post_corrupt_pptx_returns_400_when_python_pptx_available(db_ready: None) -> None:
    pytest.importorskip("pptx")
    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post(
            "/artifacts",
            files={"file": ("bad.pptx", b"not a zip pptx", _PPTX_MIME)},
        )
    assert res.status_code == 400
    assert "detail" in res.json()


@pytest.mark.asyncio
async def test_post_empty_text_pptx_returns_400_extraction_returned_no_evidence(
    db_ready: None,
) -> None:
    pytest.importorskip("pptx")
    payload = _empty_text_pptx_bytes()
    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post(
            "/artifacts",
            files={"file": ("empty.pptx", payload, _PPTX_MIME)},
        )
    assert res.status_code == 400
    assert "extraction_returned_no_evidence" in res.json()["detail"]


@pytest.mark.parametrize(
    ("modality", "upload"),
    [
        pytest.param(
            Modality.pdf,
            ("probe.pdf", _MINIMAL_HELLO_PDF, "application/pdf"),
            id="pdf",
        ),
        pytest.param(
            Modality.slide,
            ("probe.pptx", None, _PPTX_MIME),
            id="pptx",
        ),
        pytest.param(
            Modality.image,
            ("probe.png", _MINIMAL_PNG, "image/png"),
            id="image",
        ),
        pytest.param(
            Modality.audio,
            ("probe.wav", None, "audio/wav"),
            id="audio",
        ),
    ],
)
@pytest.mark.asyncio
async def test_registered_extractor_unavailable_returns_503(
    db_ready: None,
    monkeypatch: pytest.MonkeyPatch,
    modality: Modality,
    upload: tuple[str, bytes | None, str],
) -> None:
    name, payload, mime = upload
    if payload is None and name.endswith(".pptx"):
        pytest.importorskip("pptx")
        payload = _minimal_pptx_bytes()
    elif payload is None and name.endswith(".wav"):
        payload = _minimal_wav_bytes()
    assert payload is not None

    reg = ExtractorRegistry()
    reg.register(modality, _UnavailableExtractor())
    monkeypatch.setattr(artifacts_routes, "get_multimodal_extractor_registry", lambda: reg)

    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post("/artifacts", files={"file": (name, payload, mime)})
    assert res.status_code == 503


@pytest.mark.parametrize(
    "upload",
    [
        pytest.param(("x.pdf", b"%PDF-1.4 minimal", "application/pdf"), id="pdf"),
        pytest.param(("x.pptx", b"PK\x03\x04", _PPTX_MIME), id="pptx"),
        pytest.param(
            (
                "x.png",
                b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01",
                "image/png",
            ),
            id="image",
        ),
        pytest.param(("x.wav", b"RIFF\x24\x00\x00\x00WAVEfmt ", "audio/wav"), id="audio"),
    ],
)
@pytest.mark.asyncio
async def test_empty_extractor_registry_returns_400(
    db_ready: None,
    monkeypatch: pytest.MonkeyPatch,
    upload: tuple[str, bytes, str],
) -> None:
    monkeypatch.setattr(
        artifacts_routes,
        "get_multimodal_extractor_registry",
        lambda: ExtractorRegistry(),
    )
    app = create_app()
    transport = ASGITransport(app=app)
    name, body, mime = upload
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post("/artifacts", files={"file": (name, body, mime)})
    assert res.status_code == 400
    detail = res.json()["detail"].lower()
    assert (
        "extractor" in detail
        or "pdf" in detail
        or "slide" in detail
        or "image" in detail
        or "audio" in detail
    )


@pytest.mark.asyncio
async def test_video_upload_returns_400_with_video_unsupported_message(db_ready: None) -> None:
    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post(
            "/artifacts",
            files={"file": ("v.mp4", b"\x00\x00\x00", "video/mp4")},
        )
    assert res.status_code == 400
    assert "video" in res.json()["detail"].lower()


@pytest.mark.asyncio
async def test_unknown_extension_returns_400(db_ready: None) -> None:
    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post(
            "/artifacts",
            files={"file": ("x.exe", b"MZ", "application/octet-stream")},
        )
    assert res.status_code == 400


@pytest.mark.parametrize(
    ("module_path", "symbol", "filename", "body", "mime", "hint"),
    [
        pytest.param(
            "app.services.extraction.pdf_extractor",
            "PdfReader",
            "hello.pdf",
            _MINIMAL_HELLO_PDF,
            "application/pdf",
            "pdf",
            id="pdf_placeholder",
        ),
        pytest.param(
            "app.services.extraction.pptx_extractor",
            "Presentation",
            "deck.pptx",
            None,
            _PPTX_MIME,
            "pptx",
            id="pptx_placeholder",
        ),
        pytest.param(
            "app.services.extraction.image_extractor",
            "Image",
            "p.png",
            _MINIMAL_PNG,
            "image/png",
            "image",
            id="image_placeholder",
        ),
        pytest.param(
            "app.services.extraction.audio_extractor",
            "MutagenFile",
            "a.wav",
            None,
            "audio/wav",
            "audio",
            id="audio_placeholder",
        ),
    ],
)
@pytest.mark.asyncio
async def test_optional_dependency_placeholder_returns_503(
    db_ready: None,
    monkeypatch: pytest.MonkeyPatch,
    module_path: str,
    symbol: str,
    filename: str,
    body: bytes | None,
    mime: str,
    hint: str,
) -> None:
    if body is None:
        if hint == "pptx":
            pytest.importorskip("pptx")
            body = _minimal_pptx_bytes()
        else:
            body = _minimal_wav_bytes()

    mod = importlib.import_module(module_path)
    monkeypatch.setattr(mod, symbol, None)

    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post("/artifacts", files={"file": (filename, body, mime)})
    assert res.status_code == 503
    assert hint in res.json()["detail"].lower()
