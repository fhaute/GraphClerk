from __future__ import annotations

import hashlib
import os
import uuid
from io import BytesIO
from pathlib import Path

import pytest

from app.core import config as config_module
from app.models.artifact import Artifact
from app.models.enums import Modality, SourceFidelity
from app.services.errors import (
    ExtractionReturnedNoEvidenceError,
    ExtractorUnavailableError,
    PptxExtractionError,
)
from app.services.extraction.pptx_extractor import PptxExtractor
from app.services.raw_source_store import RawSourceStore


def _minimal_pptx_bytes() -> bytes:
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Slice E Title"
    body = slide.placeholders[1]
    body.text = "Body paragraph for tests."

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _empty_text_pptx_bytes() -> bytes:
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    layouts = prs.slide_layouts
    layout_idx = 6 if len(layouts) > 6 else max(0, len(layouts) - 1)
    prs.slides.add_slide(layouts[layout_idx])

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _settings(monkeypatch: pytest.MonkeyPatch, tmp_path: Path) -> object:
    monkeypatch.setenv("APP_NAME", "GraphClerk")
    monkeypatch.setenv("APP_ENV", "test")
    monkeypatch.setenv("LOG_LEVEL", "INFO")
    monkeypatch.setenv(
        "DATABASE_URL",
        os.environ.get("DATABASE_URL", "postgresql://graphclerk:graphclerk@127.0.0.1:5432/graphclerk"),
    )
    monkeypatch.setenv("QDRANT_URL", "http://localhost:6333")
    monkeypatch.setenv("ARTIFACTS_DIR", str(tmp_path / "artifacts"))
    config_module.get_settings.cache_clear()
    return config_module.get_settings()


def _artifact_for_pptx(*, raw: object, pptx_bytes: bytes) -> Artifact:
    ck = hashlib.sha256(pptx_bytes).hexdigest()
    return Artifact(
        id=uuid.uuid4(),
        filename="unit.pptx",
        title=None,
        artifact_type="pptx",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        storage_uri=raw.storage_uri,
        checksum=ck,
        size_bytes=len(pptx_bytes),
        raw_text=raw.raw_text,
        metadata_json=None,
    )


@pytest.fixture()
def python_pptx_available() -> None:
    pytest.importorskip("pptx")


def test_pptx_extractor_requires_pptx_extra(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = _settings(monkeypatch, tmp_path)
    monkeypatch.setattr("app.services.extraction.pptx_extractor.Presentation", None)
    with pytest.raises(ExtractorUnavailableError) as ei:
        PptxExtractor(settings=settings)
    assert "pptx" in str(ei.value).lower()


def test_pptx_extractor_wrong_artifact_type_raises(
    python_pptx_available: None, monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = _settings(monkeypatch, tmp_path)
    ext = PptxExtractor(settings=settings)
    bad = Artifact(
        id=uuid.uuid4(),
        filename="x.pdf",
        title=None,
        artifact_type="pdf",
        mime_type="application/pdf",
        storage_uri="x",
        checksum="a" * 64,
        size_bytes=1,
        raw_text=None,
        metadata_json=None,
    )
    with pytest.raises(PptxExtractionError):
        ext.extract(bad)


def test_pptx_extractor_extracts_title_and_body(
    python_pptx_available: None, monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = _settings(monkeypatch, tmp_path)
    pptx_bytes = _minimal_pptx_bytes()
    store = RawSourceStore(settings)
    raw = store.persist(filename="unit.pptx", content_bytes=pptx_bytes)
    artifact = _artifact_for_pptx(raw=raw, pptx_bytes=pptx_bytes)

    ext = PptxExtractor(settings=settings)
    cands = ext.extract(artifact)

    assert len(cands) >= 2
    titles = [c for c in cands if c.content_type == "slide_title"]
    bodies = [c for c in cands if c.content_type == "slide_body_text"]
    assert len(titles) >= 1
    assert len(bodies) >= 1
    t0 = titles[0]
    assert t0.modality == Modality.slide
    assert t0.source_fidelity == SourceFidelity.extracted
    assert t0.text.strip()
    assert "Slice E Title" in t0.text
    assert t0.location.get("slide_number") == 1
    assert t0.location.get("region") == "title"
    assert t0.metadata is not None
    assert t0.metadata.get("extractor") == "python-pptx"
    assert isinstance(t0.metadata.get("slide_count"), int)

    b0 = bodies[0]
    assert b0.modality == Modality.slide
    assert b0.content_type == "slide_body_text"
    assert b0.source_fidelity == SourceFidelity.extracted
    assert b0.text.strip()
    assert b0.location.get("slide_number") >= 1
    assert b0.location.get("region") == "body"


def test_pptx_extractor_bad_bytes_raises(
    python_pptx_available: None, monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = _settings(monkeypatch, tmp_path)
    bad = b"not a zip pptx"
    store = RawSourceStore(settings)
    raw = store.persist(filename="bad.pptx", content_bytes=bad)
    artifact = _artifact_for_pptx(raw=raw, pptx_bytes=bad)

    ext = PptxExtractor(settings=settings)
    with pytest.raises(PptxExtractionError):
        ext.extract(artifact)


def test_pptx_extractor_no_text_raises(
    python_pptx_available: None, monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    settings = _settings(monkeypatch, tmp_path)
    pptx_bytes = _empty_text_pptx_bytes()
    store = RawSourceStore(settings)
    raw = store.persist(filename="empty.pptx", content_bytes=pptx_bytes)
    artifact = _artifact_for_pptx(raw=raw, pptx_bytes=pptx_bytes)

    ext = PptxExtractor(settings=settings)
    with pytest.raises(ExtractionReturnedNoEvidenceError):
        ext.extract(artifact)


@pytest.mark.asyncio
async def test_post_pptx_returns_503_when_python_pptx_unavailable(
    db_ready: None, monkeypatch: pytest.MonkeyPatch
) -> None:
    import httpx
    from httpx import ASGITransport

    import app.services.extraction.pptx_extractor as pe

    monkeypatch.setattr(pe, "Presentation", None)
    from app.main import create_app

    app = create_app()
    transport = ASGITransport(app=app)
    payload = _minimal_pptx_bytes()
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post(
            "/artifacts",
            files={
                "file": (
                    "deck.pptx",
                    payload,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
    assert res.status_code == 503
    assert "pptx" in res.json()["detail"].lower()


@pytest.mark.asyncio
async def test_post_pptx_creates_artifact_and_evidence(db_ready: None) -> None:
    pytest.importorskip("pptx")
    import httpx
    from httpx import ASGITransport

    from app.main import create_app

    payload = _minimal_pptx_bytes()
    app = create_app()
    transport = ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.post(
            "/artifacts",
            files={
                "file": (
                    "deck.pptx",
                    payload,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
    assert res.status_code == 200
    data = res.json()
    assert data["artifact_type"] == "pptx"
    assert data["evidence_unit_count"] >= 1

    aid = data["artifact_id"]
    async with httpx.AsyncClient(transport=transport, base_url="http://test") as client:
        ev = await client.get(f"/artifacts/{aid}/evidence")
    assert ev.status_code == 200
    items = ev.json()["items"]
    assert len(items) >= 1
    assert any("Slice E Title" in (it.get("text") or "") for it in items)
